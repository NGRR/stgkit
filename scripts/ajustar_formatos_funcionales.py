"""Aplica formatos institucionales simples y gráficos Office nativos."""
from __future__ import annotations

import io
import shutil
import zipfile
import warnings
from copy import deepcopy
from pathlib import Path

from docx import Document
from docx.enum.style import WD_STYLE_TYPE
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Cm, Pt, RGBColor as DocxRGBColor
from docx.text.paragraph import Paragraph
from lxml import etree
from openpyxl import load_workbook
from openpyxl.chart import AreaChart, BarChart, DoughnutChart, LineChart, PieChart, RadarChart, Reference
from openpyxl.chart.label import DataLabelList
from openpyxl.chart.series import DataPoint
from PIL import Image
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt as PptxPt

ROOT = Path(__file__).resolve().parents[1]
BASE = ROOT / "base"
OUT = ROOT / "formatos" / "docx_maestro"
PPTX = ROOT / "formatos" / "pptx" / "STGND_Presentacion_Maestro_v9.pptx"
PPTX_SOURCE = ROOT / "referencias" / "legado_estaticos_v8" / "STGND_Presentacion_Riguroso_v8.pptx"
XLSX = ROOT / "formatos" / "xlsx" / "STGND_Datos_Graficos_Maestro_v9.xlsx"
REFERENCE_ZIP = BASE / "STGND - Oficio.zip"
WHITE_LOGO = ROOT / "assets" / "logo-stgnd-blanco.png"

W = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
REL = "http://schemas.openxmlformats.org/package/2006/relationships"
CT = "http://schemas.openxmlformats.org/package/2006/content-types"
WP = "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"
A = "http://schemas.openxmlformats.org/drawingml/2006/main"
C = "http://schemas.openxmlformats.org/drawingml/2006/chart"
NS = {"w": W, "r": R, "rel": REL, "ct": CT, "wp": WP, "a": A, "c": C}


def set_paragraph(paragraph, text: str) -> None:
    if paragraph.runs:
        paragraph.runs[0].text = text
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.add_run(text)


def remove_paragraph(paragraph) -> None:
    element = paragraph._element
    element.getparent().remove(element)


def remove_table(table) -> None:
    element = table._element
    element.getparent().remove(element)


def normalize_normal_style(document: Document) -> None:
    style = document.styles["Normal"]
    style.font.name = "Aptos"
    style.font.size = Pt(11)


def ensure_semantic_styles(document: Document) -> None:
    """Crea estilos interoperables para TOC, navegación y contenido automático."""
    normalize_normal_style(document)
    specs = {
        "Title": (26, "Aptos Display", "4C2B46", None),
        "Subtitle": (13, "Aptos", "666666", None),
        "Heading 1": (18, "Aptos Display", "4C2B46", 0),
        "Heading 2": (14, "Aptos Display", "E97700", 1),
        "Heading 3": (11, "Aptos", "4C2B46", 2),
        "Caption": (9, "Aptos", "666666", None),
    }
    for name, (size, font, color, outline) in specs.items():
        try:
            style = document.styles[name]
        except KeyError:
            style = document.styles.add_style(name, WD_STYLE_TYPE.PARAGRAPH)
        style.base_style = document.styles["Normal"]
        style.font.name = font
        style.font.size = Pt(size)
        style.font.color.rgb = DocxRGBColor.from_string(color)
        style.font.bold = name in {"Title", "Heading 1", "Heading 2", "Heading 3"}
        style.paragraph_format.space_before = Pt(12 if name.startswith("Heading") else 0)
        style.paragraph_format.space_after = Pt(6)
        properties = style._element.get_or_add_pPr()
        for child in list(properties):
            if child.tag == f"{{{W}}}outlineLvl":
                properties.remove(child)
        if outline is not None:
            level = etree.SubElement(properties, f"{{{W}}}outlineLvl")
            level.set(f"{{{W}}}val", str(outline))
            etree.SubElement(properties, f"{{{W}}}keepNext")
        if style._element.find(f"{{{W}}}qFormat") is None:
            etree.SubElement(style._element, f"{{{W}}}qFormat")


def clean_master(master: str) -> Document:
    document = Document(BASE / master)
    body = document._element.body
    for child in list(body):
        if child.tag != f"{{{W}}}sectPr":
            body.remove(child)
    ensure_semantic_styles(document)
    return document


def style_table(table, *, header: bool = True) -> None:
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.autofit = True
    borders = table._tbl.tblPr.find(f"{{{W}}}tblBorders")
    if borders is None:
        borders = etree.SubElement(table._tbl.tblPr, f"{{{W}}}tblBorders")
    for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
        border = borders.find(f"{{{W}}}{edge}")
        if border is None:
            border = etree.SubElement(borders, f"{{{W}}}{edge}")
        border.set(f"{{{W}}}val", "single")
        border.set(f"{{{W}}}sz", "4")
        border.set(f"{{{W}}}color", "D9D5D8")
    for row_index, row in enumerate(table.rows):
        for cell in row.cells:
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            if header and row_index == 0:
                shading = cell._tc.get_or_add_tcPr().find(f"{{{W}}}shd")
                if shading is None:
                    shading = etree.SubElement(cell._tc.get_or_add_tcPr(), f"{{{W}}}shd")
                shading.set(f"{{{W}}}fill", "4C2B46")
            for paragraph in cell.paragraphs:
                paragraph.style = "Normal"
                for run in paragraph.runs:
                    run.font.name = "Aptos"
                    run.font.size = Pt(11)
                    if header and row_index == 0:
                        run.font.bold = True
                        run.font.color.rgb = DocxRGBColor(255, 255, 255)
def add_metadata_table(document: Document, rows: list[tuple[str, str]]) -> None:
    table = document.add_table(rows=len(rows), cols=2)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for index, (label, value) in enumerate(rows):
        table.cell(index, 0).text = label
        table.cell(index, 1).text = value
        shade = etree.SubElement(table.cell(index, 0)._tc.get_or_add_tcPr(), f"{{{W}}}shd")
        shade.set(f"{{{W}}}fill", "EFEDEF")
        for run in table.cell(index, 0).paragraphs[0].runs:
            run.font.bold = True
            run.font.color.rgb = DocxRGBColor.from_string("4C2B46")
    style_table(table, header=False)


def add_dynamic_table(document: Document, title: str, headers: list[str]) -> None:
    heading = document.add_paragraph(title)
    heading.style = document.styles["Heading 1"]
    append_dynamic_table(document, headers)


def set_core(document: Document, title: str) -> None:
    document.core_properties.title = title
    document.core_properties.subject = "Formato institucional STGND"


def save_atomic(document: Document, target: Path) -> None:
    temporary = target.with_suffix(".docx.tmp")
    document.save(temporary)
    temporary.replace(target)


def reference_document(name: str) -> Document:
    with zipfile.ZipFile(REFERENCE_ZIP) as package:
        return Document(io.BytesIO(package.read(name)))


def append_dynamic_table(document: Document, headers: list[str]) -> None:
    table = document.add_table(rows=2, cols=len(headers))
    table.autofit = True
    borders = etree.SubElement(table._tbl.tblPr, f"{{{W}}}tblBorders")
    for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
        border = etree.SubElement(borders, f"{{{W}}}{edge}")
        border.set(f"{{{W}}}val", "single")
        border.set(f"{{{W}}}sz", "4")
        border.set(f"{{{W}}}color", "D9D5D8")
    for index, header in enumerate(headers):
        cell = table.rows[0].cells[index]
        cell.text = header
        shading = etree.SubElement(cell._tc.get_or_add_tcPr(), f"{{{W}}}shd")
        shading.set(f"{{{W}}}fill", "4C2B46")
        for run in cell.paragraphs[0].runs:
            run.font.name = "Aptos"
            run.font.size = Pt(11)
            run.font.bold = True
            run.font.color.rgb = DocxRGBColor(255, 255, 255)
    for index, cell in enumerate(table.rows[1].cells):
        cell.text = "{{dynamic_rows}}" if index == 0 else ""
        for run in cell.paragraphs[0].runs:
            run.font.name = "Aptos"
            run.font.size = Pt(11)


def build_oficio() -> None:
    document = reference_document("STGND - Oficio.docx")
    values = {
        0: "A             :     \t{{destinatario}}",
        1: "{{cargo_destinatario}}",
        2: "{{institucion_destinataria}}",
        3: "DE          :     \t{{remitente}}",
        4: "{{cargo_remitente}}",
        5: "CORTE SUPREMA DE JUSTICIA",
        7: "{{cuerpo}}",
        9: "Sin otro particular,",
        20: "Distribución:",
        21: "{{distribucion}} · Archivo Secretaría Técnica de Igualdad de Género y No Discriminación",
    }
    for index, text in values.items():
        set_paragraph(document.paragraphs[index], text)
    cell = document.tables[0].cell(0, 1)
    set_paragraph(cell.paragraphs[0], "OFICIO N.° {{numero_documento}}")
    set_paragraph(cell.paragraphs[1], "MAT.:\t{{materia}}")
    set_paragraph(cell.paragraphs[4], "REF.:\t{{antecedente}}")
    set_paragraph(cell.paragraphs[8], "{{ciudad_upper}}, {{fecha}}")
    normalize_normal_style(document)
    set_core(document, "Oficio institucional")
    save_atomic(document, OUT / "STGND_Oficio_Maestro_v9.docx")


def build_circular() -> None:
    document = reference_document("STGND - Oficio circular.docx")
    values = {
        0: "A             :     \t{{destinatarios}}",
        1: "",
        3: "DE          :     \t{{remitente}}",
        4: "{{cargo_remitente}}",
        5: "CORTE SUPREMA DE JUSTICIA",
        7: "{{cuerpo}}",
        9: "",
        11: "",
        12: "Agradeciendo su atención,",
        20: "Distribución:",
        21: "{{distribucion}}",
        22: "",
        23: "Archivo Secretaría Técnica de Igualdad de Género y No Discriminación",
    }
    for index, text in values.items():
        set_paragraph(document.paragraphs[index], text)
    cell = document.tables[0].cell(0, 1)
    set_paragraph(cell.paragraphs[0], "OFICIO CIRCULAR N.° {{numero_documento}}")
    set_paragraph(cell.paragraphs[3], "MAT.:\t{{materia}}")
    set_paragraph(cell.paragraphs[5], "REF.:\t{{antecedente}}")
    set_paragraph(cell.paragraphs[6], "")
    set_paragraph(cell.paragraphs[9], "{{ciudad_upper}}, {{fecha}}")
    normalize_normal_style(document)
    set_core(document, "Oficio circular institucional")
    save_atomic(document, OUT / "STGND_Oficio_Circular_Maestro_v9.docx")


def build_agenda() -> None:
    document = clean_master("DocumentoMaestro.docx")
    title = document.add_paragraph("AGENDA DE REUNIÓN · {{titulo}}")
    title.style = document.styles["Title"]
    subtitle = document.add_paragraph("Planificación operativa de sesión, sin portada ni secciones editoriales.")
    subtitle.style = document.styles["Subtitle"]
    add_metadata_table(document, [
        ("Sesión", "N.° {{numero_sesion}}"),
        ("Fecha", "{{ciudad}}, {{fecha}}"),
        ("Horario", "{{hora_inicio}} a {{hora_termino}}"),
        ("Modalidad", "{{modalidad}}"),
        ("Lugar", "{{lugar}}"),
    ])
    heading = document.add_paragraph("Objetivo de la reunión")
    heading.style = document.styles["Heading 1"]
    document.add_paragraph("{{objetivo}}", style="Normal")
    add_dynamic_table(document, "Agenda de trabajo", ["Hora", "Tema / actividad", "Responsable", "Duración"])
    heading = document.add_paragraph("Documentos preparatorios")
    heading.style = document.styles["Heading 1"]
    document.add_paragraph("{{documentos_preparatorios}}", style="Normal")
    set_core(document, "Agenda de reunión")
    save_atomic(document, OUT / "STGND_Agenda_Reunion_Maestro_v9.docx")


def build_acta() -> None:
    document = reference_document("STGND - Acta.docx")
    values = {
        0: "ACTA {{titulo_upper}}",
        1: "Sesión N.° {{numero_sesion}} · {{ciudad}}, {{fecha}}",
        2: "Horario: {{hora_inicio}} a {{hora_termino}}",
        3: "Modalidad: {{modalidad}}",
        4: "Lugar: {{lugar}}",
        5: "Participantes: {{participantes}}",
        7: "OBJETIVO DE LA REUNIÓN: {{objetivo}}",
        9: "TEMAS TRATADOS",
        10: "{{desarrollo}}",
        11: "",
        12: "",
        13: "",
        15: "DOCUMENTOS PREPARATORIOS",
        16: "{{documentos_preparatorios}}",
        17: "",
        18: "",
        23: "ACUERDOS Y SEGUIMIENTO",
        24: "Registro de acuerdos, responsables, plazos y estado.",
        25: "",
        26: "",
        27: "PRÓXIMA REUNIÓN",
        28: "{{proxima_reunion}}",
    }
    for index, text in values.items():
        set_paragraph(document.paragraphs[index], text)
    remove_table(document.tables[0])
    append_dynamic_table(document, ["Acuerdo", "Responsable", "Plazo", "Estado"])
    normalize_normal_style(document)
    set_core(document, "Acta de reunión")
    save_atomic(document, OUT / "STGND_Acta_Reunion_Maestro_v9.docx")


def build_minuta() -> None:
    document = clean_master("DocumentoMaestro2.docx")
    title = document.add_paragraph("MINUTA DE ACTIVIDAD · {{titulo}}")
    title.style = document.styles["Title"]
    subtitle = document.add_paragraph("Registro breve de coordinación, contenido y seguimiento.")
    subtitle.style = document.styles["Subtitle"]
    add_metadata_table(document, [
        ("Fecha", "{{fecha}}"),
        ("Horario", "{{hora_inicio}} a {{hora_termino}}"),
        ("Lugar / modalidad", "{{lugar}} · {{modalidad}}"),
        ("Organiza", "{{responsable}}"),
        ("Participantes", "{{participantes}}"),
        ("Difusión", "{{difusion}}"),
    ])
    for heading_text, content in (
        ("Objetivo", "{{objetivo}}"),
        ("Síntesis", "{{resumen}}"),
        ("Desarrollo", "{{cuerpo}}"),
        ("Justificación", "{{idea_clave}}"),
    ):
        heading = document.add_paragraph(heading_text)
        heading.style = document.styles["Heading 1"]
        document.add_paragraph(content, style="Normal")
    add_dynamic_table(document, "Acuerdos y seguimiento", ["Acuerdo", "Responsable", "Plazo", "Estado"])
    set_core(document, "Minuta de actividad")
    save_atomic(document, OUT / "STGND_Minuta_Actividad_Maestro_v9.docx")


def build_programa() -> None:
    document = Document(BASE / "ProgramaHito_lanzamiento.docx")
    set_paragraph(document.paragraphs[1], "PROGRAMA")
    set_paragraph(document.paragraphs[2], "{{titulo_upper}}")
    set_paragraph(document.paragraphs[3], "{{descripcion_actividad_upper}}")
    set_paragraph(document.paragraphs[4], "{{fecha}}")
    set_paragraph(document.paragraphs[5], "{{hora_inicio}} · {{lugar}} · {{modalidad}}")
    table = document.tables[0]
    for row in list(table.rows)[2:]:
        table._tbl.remove(row._tr)
    for index, cell in enumerate(table.rows[1].cells):
        set_paragraph(cell.paragraphs[0], "{{dynamic_rows}}" if index == 0 else "")
    normalize_normal_style(document)
    set_core(document, "Programa de actividad")
    save_atomic(document, OUT / "STGND_Programa_Actividad_Maestro_v9.docx")


def set_xml_paragraph(paragraph: etree._Element, text: str) -> None:
    nodes = paragraph.xpath(".//w:t", namespaces=NS)
    if not nodes:
        return
    nodes[0].text = text
    nodes[0].set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
    for node in nodes[1:]:
        node.text = ""


def build_invitation() -> None:
    source = BASE / "Invitacion_Hito_Guia.docx"
    with zipfile.ZipFile(source) as package:
        entries = [(info, package.read(info.filename)) for info in package.infolist()]
    root = etree.fromstring(dict((i.filename, d) for i, d in entries)["word/document.xml"])
    paragraphs = root.xpath("//w:p", namespaces=NS)
    replacements = {
        2: "{{preheader}}",
        3: "{{titulo}}",
        4: "Invitación institucional",
        7: "{{cuerpo_invitacion}}",
        9: "{{fecha}} · {{hora_inicio}} · {{lugar}} · {{modalidad}}",
        12: "{{ciudad}}, {{fecha}}",
        14: "{{cta}} · {{url_cta}}",
    }
    for index, text in replacements.items():
        set_xml_paragraph(paragraphs[index], text)
    document_xml = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
    target = OUT / "STGND_Invitacion_Email_Maestro_v9.docx"
    temporary = target.with_suffix(".docx.tmp")
    with zipfile.ZipFile(temporary, "w") as output:
        for info, data in entries:
            output.writestr(info, document_xml if info.filename == "word/document.xml" else data)
    temporary.replace(target)
    document = Document(target)
    normalize_normal_style(document)
    set_core(document, "Invitación institucional por correo")
    save_atomic(document, target)


def signature_table() -> etree._Element:
    xml = f'''<w:tbl xmlns:w="{W}">
      <w:tblPr><w:tblW w:w="4200" w:type="dxa"/><w:jc w:val="center"/><w:tblBorders>
      <w:top w:val="single" w:sz="6" w:color="D9D5D8"/><w:left w:val="single" w:sz="6" w:color="D9D5D8"/>
      <w:bottom w:val="single" w:sz="6" w:color="D9D5D8"/><w:right w:val="single" w:sz="6" w:color="D9D5D8"/>
      </w:tblBorders></w:tblPr><w:tblGrid><w:gridCol w:w="4200"/></w:tblGrid><w:tr><w:trPr><w:trHeight w:val="900"/></w:trPr><w:tc>
      <w:tcPr><w:tcW w:w="4200" w:type="dxa"/><w:vAlign w:val="center"/><w:shd w:fill="F7F5F6"/></w:tcPr>
      <w:p><w:pPr><w:jc w:val="center"/><w:spacing w:after="80"/></w:pPr><w:r><w:rPr><w:rFonts w:ascii="Aptos" w:hAnsi="Aptos"/><w:b/><w:color w:val="4C2B46"/><w:sz w:val="22"/></w:rPr><w:t>FIRMA ELECTRÓNICA</w:t></w:r></w:p>
      <w:p><w:pPr><w:jc w:val="center"/></w:pPr><w:r><w:rPr><w:vanish/></w:rPr><w:t>FIRMADIGITAL</w:t></w:r><w:r><w:rPr><w:rFonts w:ascii="Aptos" w:hAnsi="Aptos"/><w:color w:val="666666"/><w:sz w:val="18"/></w:rPr><w:t>Espacio reservado para validación institucional</w:t></w:r></w:p>
      </w:tc></w:tr></w:tbl>'''
    return etree.fromstring(xml.encode("utf-8"))


def normalize_signature_box(path: Path, *, ensure: bool) -> None:
    with zipfile.ZipFile(path) as package:
        entries = [(info, package.read(info.filename)) for info in package.infolist()]
    entry_map = dict((info.filename, data) for info, data in entries)
    root = etree.fromstring(entry_map["word/document.xml"])
    body = root.find(f".//{{{W}}}body")
    removed = False
    for anchor in list(root.xpath("//wp:anchor", namespaces=NS)):
        text = "".join(anchor.xpath(".//w:t/text()", namespaces=NS))
        if "FIRMADIGITAL" not in text:
            continue
        run = anchor
        while run is not None and run.tag != f"{{{W}}}r":
            run = run.getparent()
        if run is not None and run.getparent() is not None:
            run.getparent().remove(run)
            removed = True
    if ensure:
        anchor_node = next((node for node in body if "Distribución" in "".join(node.xpath(".//w:t/text()", namespaces=NS))), None)
        if anchor_node is None:
            anchor_node = body.find(f"{{{W}}}sectPr")
        body.insert(body.index(anchor_node), signature_table())
    replacement = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
    temporary = path.with_suffix(".docx.tmp")
    with zipfile.ZipFile(temporary, "w") as output:
        for info, data in entries:
            output.writestr(info, replacement if info.filename == "word/document.xml" else data)
    temporary.replace(path)


def enforce_formal_black(path: Path) -> None:
    """Mantiene la correspondencia oficial en negro, sin estilos editoriales."""
    with zipfile.ZipFile(path) as package:
        entries = [(info, package.read(info.filename)) for info in package.infolist()]
    root = etree.fromstring(dict((info.filename, data) for info, data in entries)["word/document.xml"])
    for run in root.xpath("//w:r[.//w:t]", namespaces=NS):
        properties = run.find(f"{{{W}}}rPr")
        if properties is None:
            properties = etree.Element(f"{{{W}}}rPr")
            run.insert(0, properties)
        color = properties.find(f"{{{W}}}color")
        if color is None:
            color = etree.SubElement(properties, f"{{{W}}}color")
        color.set(f"{{{W}}}val", "000000")
        for attribute in ("themeColor", "themeTint", "themeShade"):
            color.attrib.pop(f"{{{W}}}{attribute}", None)
    replacement = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
    temporary = path.with_suffix(".docx.tmp")
    with zipfile.ZipFile(temporary, "w") as output:
        for info, data in entries:
            output.writestr(info, replacement if info.filename == "word/document.xml" else data)
    temporary.replace(path)


def apply_semantic_styles(path: Path) -> None:
    document = Document(path)
    ensure_semantic_styles(document)
    heading_terms = {
        "Resumen", "Antecedentes", "Desarrollo", "Síntesis", "Objetivo", "Justificación",
        "Control documental", "Control normativo", "Ficha de datos", "Control de tareas",
        "Agenda de trabajo", "Documentos preparatorios", "Acuerdos y seguimiento",
        "PRÓXIMA REUNIÓN", "ACUERDOS Y SEGUIMIENTO", "TEMAS TRATADOS",
    }
    title_assigned = False
    # document.paragraphs omite texto dentro de tablas y cuadros de texto.
    # Recorrer el árbol completo asigna semántica también a oficios e invitaciones.
    paragraphs = [Paragraph(element, document) for element in document.element.body.iter(f"{{{W}}}p")]
    for paragraph in paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        if not title_assigned and ("{{titulo" in text or text.startswith(("ACTA", "AGENDA", "MINUTA", "PROGRAMA", "OFICIO"))):
            paragraph.style = document.styles["Title"]
            title_assigned = True
            continue
        if text in heading_terms or text.upper() in heading_terms:
            paragraph.style = document.styles["Heading 1"]
        elif text.startswith(("FIGURA", "TABLA", "ESQUEMA")):
            paragraph.style = document.styles["Caption"]
        elif "{{" in text and paragraph.style.name not in {"Title", "Heading 1", "Heading 2", "Heading 3"}:
            paragraph.style = document.styles["Normal"]
    save_atomic(document, path)


CHART_LIBRARY = [
    ("Columnas agrupadas", XL_CHART_TYPE.COLUMN_CLUSTERED),
    ("Columnas apiladas", XL_CHART_TYPE.COLUMN_STACKED),
    ("Columnas apiladas 100 %", XL_CHART_TYPE.COLUMN_STACKED_100),
    ("Barras horizontales", XL_CHART_TYPE.BAR_CLUSTERED),
    ("Barras apiladas", XL_CHART_TYPE.BAR_STACKED),
    ("Gráfico de torta", XL_CHART_TYPE.PIE),
    ("Gráfico de anillo", XL_CHART_TYPE.DOUGHNUT),
    ("Evolución temporal", XL_CHART_TYPE.LINE_MARKERS),
    ("Área acumulada", XL_CHART_TYPE.AREA),
    ("Perfil comparado", XL_CHART_TYPE.RADAR),
]


def chart_carrier(chart_type) -> tuple[bytes, bytes, bytes]:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    data = ChartData()
    data.categories = ["Categoría A", "Categoría B", "Categoría C", "Categoría D", "Categoría E"]
    data.add_series("Resultado", (82, 64, 47, 31, 18))
    if chart_type not in {XL_CHART_TYPE.PIE, XL_CHART_TYPE.DOUGHNUT}:
        data.add_series("Referencia", (70, 58, 52, 40, 25))
    chart = slide.shapes.add_chart(chart_type, Inches(1), Inches(1), Inches(7), Inches(4), data).chart
    chart.has_legend = len(chart.series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    try:
        chart.value_axis.maximum_scale = 100
        chart.value_axis.minimum_scale = 0
        chart.value_axis.has_major_gridlines = False
    except ValueError:
        pass
    palette = [RGBColor(76, 43, 70), RGBColor(233, 119, 0), RGBColor(137, 121, 132), RGBColor(209, 199, 205), RGBColor(98, 67, 91)]
    for index, series in enumerate(chart.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = palette[index % len(palette)]
    if chart_type in {XL_CHART_TYPE.PIE, XL_CHART_TYPE.DOUGHNUT}:
        for point, color in zip(chart.series[0].points, palette):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = color
    chart.series[0].has_data_labels = True
    stream = io.BytesIO()
    presentation.save(stream)
    with zipfile.ZipFile(io.BytesIO(stream.getvalue())) as package:
        chart_xml = package.read("ppt/charts/chart1.xml")
        chart_rels = package.read("ppt/charts/_rels/chart1.xml.rels")
        workbook = package.read("ppt/embeddings/Microsoft_Excel_Sheet1.xlsx")
    return chart_xml, chart_rels, workbook


def chart_paragraph(relationship_id: str, identifier: int) -> etree._Element:
    xml = f'''<w:p xmlns:w="{W}" xmlns:r="{R}" xmlns:wp="{WP}" xmlns:a="{A}" xmlns:c="{C}">
      <w:pPr><w:jc w:val="center"/></w:pPr><w:r><w:drawing><wp:inline distT="0" distB="0" distL="0" distR="0">
      <wp:extent cx="5486400" cy="3200400"/><wp:effectExtent l="0" t="0" r="0" b="0"/><wp:docPr id="{990 + identifier}" name="Gráfico Office nativo {identifier}"/>
      <wp:cNvGraphicFramePr/><a:graphic><a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/chart"><c:chart r:id="{relationship_id}"/></a:graphicData></a:graphic>
      </wp:inline></w:drawing></w:r></w:p>'''
    return etree.fromstring(xml.encode("utf-8"))


def graphic_heading(title: str, index: int) -> etree._Element:
    xml = f'''<w:p xmlns:w="{W}"><w:pPr><w:pStyle w:val="Heading2"/><w:keepNext/><w:pageBreakBefore/></w:pPr>
    <w:r><w:rPr><w:rFonts w:ascii="Aptos Display" w:hAnsi="Aptos Display"/><w:b/><w:color w:val="4C2B46"/><w:sz w:val="28"/></w:rPr>
    <w:t>FIGURA {index:02d} · {title}</w:t></w:r></w:p>'''
    return etree.fromstring(xml.encode("utf-8"))


def inject_native_charts(path: Path, chart_specs: list[tuple[str, object]]) -> None:
    with zipfile.ZipFile(path) as package:
        entries = [(info, package.read(info.filename)) for info in package.infolist()]
    entry_map = dict((info.filename, data) for info, data in entries)
    document = etree.fromstring(entry_map["word/document.xml"])
    body = document.find(f".//{{{W}}}body")
    children = list(body)
    start = next((i for i, node in enumerate(children) if "58%PARTICIPACIÓN" in "".join(node.xpath(".//w:t/text()", namespaces=NS))), None)
    end = next((i for i, node in enumerate(children) if "FIGURA 01Distribución comparada" in "".join(node.xpath(".//w:t/text()", namespaces=NS))), None)
    if start is None or end is None:
        raise ValueError(f"No se encontró el gráfico tabular en {path.name}")
    for node in children[start : end + 1]:
        body.remove(node)

    relationships = etree.fromstring(entry_map["word/_rels/document.xml.rels"])
    types = etree.fromstring(entry_map["[Content_Types].xml"])
    existing = {node.get("PartName") for node in types}
    additions: dict[str, bytes] = {}
    insertion = min(start, len(body) - 1)
    used = {node.get("Id") for node in relationships}
    for index, (title, chart_type) in enumerate(chart_specs, 1):
        chart_xml, chart_rels, workbook = chart_carrier(chart_type)
        relationship_id = f"rIdChart{index}"
        while relationship_id in used:
            relationship_id += "x"
        used.add(relationship_id)
        relation = etree.SubElement(relationships, f"{{{REL}}}Relationship")
        relation.set("Id", relationship_id)
        relation.set("Type", "http://schemas.openxmlformats.org/officeDocument/2006/relationships/chart")
        relation.set("Target", f"charts/chart{index}.xml")
        body.insert(insertion, graphic_heading(title, index))
        body.insert(insertion + 1, chart_paragraph(relationship_id, index))
        insertion += 2
        workbook_name = f"Microsoft_Excel_Sheet{index}.xlsx"
        additions[f"word/charts/chart{index}.xml"] = chart_xml
        additions[f"word/charts/_rels/chart{index}.xml.rels"] = chart_rels.replace(b"Microsoft_Excel_Sheet1.xlsx", workbook_name.encode("ascii"))
        additions[f"word/embeddings/{workbook_name}"] = workbook
        for part_name, content_type in (
            (f"/word/charts/chart{index}.xml", "application/vnd.openxmlformats-officedocument.drawingml.chart+xml"),
            (f"/word/embeddings/{workbook_name}", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
        ):
            if part_name not in existing:
                override = etree.SubElement(types, f"{{{CT}}}Override")
                override.set("PartName", part_name)
                override.set("ContentType", content_type)
                existing.add(part_name)
    replacements = {
        "word/document.xml": etree.tostring(document, xml_declaration=True, encoding="UTF-8", standalone=True),
        "word/_rels/document.xml.rels": etree.tostring(relationships, xml_declaration=True, encoding="UTF-8", standalone=True),
        "[Content_Types].xml": etree.tostring(types, xml_declaration=True, encoding="UTF-8", standalone=True),
    }
    temporary = path.with_suffix(".docx.tmp")
    with zipfile.ZipFile(temporary, "w") as output:
        for info, data in entries:
            if info.filename not in additions:
                output.writestr(info, replacements.get(info.filename, data))
        for name, data in additions.items():
            output.writestr(name, data)
    temporary.replace(path)


def add_structured_data_examples(path: Path) -> None:
    document = Document(path)
    ensure_semantic_styles(document)
    heading = document.add_paragraph("TABLA 01 · Comparación de indicadores")
    heading.style = document.styles["Heading 2"]
    table = document.add_table(rows=4, cols=4)
    values = [
        ["Indicador", "Línea base", "Resultado", "Variación"],
        ["Participación", "42 %", "58 %", "+16 pp"],
        ["Casos analizados", "18", "23", "+5"],
        ["Índice comparado", "3,8", "4,2", "+0,4"],
    ]
    for row, row_values in zip(table.rows, values):
        for cell, value in zip(row.cells, row_values):
            cell.text = value
    style_table(table)
    heading = document.add_paragraph("ESQUEMA 01 · Flujo de decisión")
    heading.style = document.styles["Heading 2"]
    flow = document.add_table(rows=1, cols=7)
    for index, value in enumerate(["Entrada", "→", "Validación", "→", "Análisis", "→", "Decisión"]):
        cell = flow.cell(0, index)
        cell.text = value
        if value != "→":
            shade = etree.SubElement(cell._tc.get_or_add_tcPr(), f"{{{W}}}shd")
            shade.set(f"{{{W}}}fill", "4C2B46" if index in {0, 4} else "E97700")
            for run in cell.paragraphs[0].runs:
                run.font.bold = True
                run.font.color.rgb = DocxRGBColor(255, 255, 255)
        cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
    style_table(flow, header=False)
    document.add_paragraph("Los gráficos son objetos Office editables; las tablas y el esquema usan estilos institucionales reutilizables.", style="Caption")
    save_atomic(document, path)


def slide_text(slide) -> str:
    return " ".join(shape.text for shape in slide.shapes if hasattr(shape, "text_frame") and shape.has_text_frame)


def normalize_unique_zip(path: Path) -> None:
    """Reescribe el paquete conservando la última versión de cada parte OOXML."""
    with zipfile.ZipFile(path) as package:
        entries = {info.filename: (info, package.read(info)) for info in package.infolist()}
    temporary = path.with_suffix(path.suffix + ".unique")
    with zipfile.ZipFile(temporary, "w") as output:
        for info, data in entries.values():
            output.writestr(info, data)
    temporary.replace(path)


def remove_slide(presentation: Presentation, slide) -> None:
    slide_id = slide.slide_id
    for node in list(presentation.slides._sldIdLst):
        if int(node.id) == slide_id:
            relationship_id = node.rId
            presentation.part.drop_rel(relationship_id)
            presentation.slides._sldIdLst.remove(node)
            break


def set_slide_background(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(color)


def ensure_white_logo() -> Path:
    """Crea una versión blanca con transparencia real para máxima compatibilidad."""
    source = Image.open(ROOT / "assets" / "logo-stgnd.png").convert("RGBA")
    white = Image.new("RGBA", source.size, (255, 255, 255, 0))
    white.putalpha(source.getchannel("A"))
    white.save(WHITE_LOGO, optimize=True)
    return WHITE_LOGO


def add_textbox(slide, text: str, x, y, width, height, *, size: int, color: str, bold: bool = False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(x, y, width, height)
    frame = shape.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    run = paragraph.runs[0]
    run.font.name = "Aptos Display" if size >= 20 else "Aptos"
    run.font.size = PptxPt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)
    return shape


def add_cover(presentation: Presentation, *, purple: bool):
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    background = "4C2B46" if purple else "FFFFFF"
    primary = "FFFFFF" if purple else "4C2B46"
    set_slide_background(slide, background)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(0.16), Inches(4.8))
    accent.fill.solid(); accent.fill.fore_color.rgb = RGBColor.from_string("E97700"); accent.line.fill.background()
    logo_path = ensure_white_logo() if purple else ROOT / "assets/logo-stgnd.png"
    slide.shapes.add_picture(str(logo_path), Inches(1.2), Inches(0.65), width=Inches(4.9))
    add_textbox(slide, "PRESENTACIÓN INSTITUCIONAL", Inches(1.2), Inches(3.0), Inches(8.8), Inches(0.45), size=12, color="E97700", bold=True)
    add_textbox(slide, "Título de la presentación", Inches(1.2), Inches(3.45), Inches(9.6), Inches(1.0), size=34, color=primary, bold=True)
    add_textbox(slide, "Subtítulo, período y unidad responsable", Inches(1.2), Inches(4.45), Inches(8.8), Inches(0.55), size=16, color=primary)
    add_textbox(slide, f"OPCIÓN DE PORTADA · {'MORADA' if purple else 'BLANCA'}", Inches(1.2), Inches(6.35), Inches(5), Inches(0.3), size=8, color="E97700")
    return slide


def add_chart_slide(presentation: Presentation, title: str, chart_type) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    set_slide_background(slide, "FFFFFF")
    add_textbox(slide, title, Inches(0.8), Inches(0.55), Inches(10.5), Inches(0.55), size=25, color="4C2B46", bold=True)
    add_textbox(slide, "EJEMPLO NATIVO · gráfico editable en PowerPoint", Inches(0.8), Inches(1.15), Inches(7.5), Inches(0.3), size=9, color="E97700", bold=True)
    data = ChartData(); data.categories = ["A", "B", "C", "D", "E"]; data.add_series("Resultado", (82, 64, 47, 31, 18))
    if chart_type not in {XL_CHART_TYPE.PIE, XL_CHART_TYPE.DOUGHNUT}:
        data.add_series("Referencia", (70, 58, 52, 40, 25))
    chart = slide.shapes.add_chart(chart_type, Inches(1.0), Inches(1.7), Inches(9.8), Inches(4.7), data).chart
    chart.has_title = False; chart.has_legend = len(chart.series) > 1
    if chart.has_legend: chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    colors = [RGBColor.from_string("4C2B46"), RGBColor.from_string("E97700")]
    for series, color in zip(chart.series, colors):
        series.format.fill.solid(); series.format.fill.fore_color.rgb = color
    if chart_type in {XL_CHART_TYPE.PIE, XL_CHART_TYPE.DOUGHNUT}:
        palette = ["4C2B46", "E97700", "896F83", "CFC6CC", "623F5B"]
        for point, color in zip(chart.series[0].points, palette):
            point.format.fill.solid(); point.format.fill.fore_color.rgb = RGBColor.from_string(color)


def fix_presentation() -> None:
    # Partir siempre de la referencia evita acumulación de diapositivas y
    # colisiones de nombres internos en reconstrucciones sucesivas.
    shutil.copy2(PPTX_SOURCE, PPTX)
    normalize_unique_zip(PPTX)
    presentation = Presentation(PPTX)
    for slide in list(presentation.slides):
        if "OPCIÓN DE PORTADA" in slide_text(slide) or "EJEMPLO NATIVO" in slide_text(slide):
            remove_slide(presentation, slide)
    for slide in presentation.slides:
        for shape in list(slide.shapes):
            rotation = float(shape.rotation or 0)
            if abs(rotation - 270) < 0.1 or abs(rotation - 90) < 0.1:
                slide.shapes._spTree.remove(shape._element)

    chart_slide = next((slide for slide in presentation.slides if "Distribución comparada" in slide_text(slide)), None)
    if chart_slide is not None:
        for shape in list(chart_slide.shapes):
            x, y = shape.left / 914400, shape.top / 914400
            if 0.8 <= x <= 5.9 and 3.75 <= y <= 5.75:
                chart_slide.shapes._spTree.remove(shape._element)
        data = ChartData(); data.categories = ["A", "B", "C", "D", "E"]; data.add_series("Resultado", (82, 64, 47, 31, 18))
        chart = chart_slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.95), Inches(3.75), Inches(5.15), Inches(2.05), data).chart
        chart.has_legend = False; chart.series[0].format.fill.solid(); chart.series[0].format.fill.fore_color.rgb = RGBColor.from_string("4C2B46")

    closing = presentation.slides[-1]
    for shape in list(closing.shapes):
        closing.shapes._spTree.remove(shape._element)
    set_slide_background(closing, "4C2B46")
    closing.shapes.add_picture(str(ensure_white_logo()), Inches(3.6), Inches(2.35), width=Inches(5.1))
    add_textbox(closing, "Gracias", Inches(1.0), Inches(4.65), Inches(10.0), Inches(0.7), size=30, color="FFFFFF", bold=True, align=PP_ALIGN.CENTER)
    add_textbox(closing, "secretariadegenero.pjud.cl", Inches(1.0), Inches(5.35), Inches(10.0), Inches(0.35), size=11, color="FFFFFF", align=PP_ALIGN.CENTER)

    purple_cover = add_cover(presentation, purple=True)
    white_cover = add_cover(presentation, purple=False)
    add_chart_slide(presentation, "Distribución porcentual", XL_CHART_TYPE.PIE)
    add_chart_slide(presentation, "Evolución temporal", XL_CHART_TYPE.LINE_MARKERS)
    add_chart_slide(presentation, "Comparación de resultados", XL_CHART_TYPE.COLUMN_CLUSTERED)
    slide_ids = presentation.slides._sldIdLst
    purple_id, white_id = slide_ids[-5], slide_ids[-4]
    slide_ids.remove(purple_id); slide_ids.remove(white_id)
    slide_ids.insert(0, white_id); slide_ids.insert(0, purple_id)
    closing_id = next(node for node in slide_ids if int(node.id) == closing.slide_id)
    slide_ids.remove(closing_id); slide_ids.append(closing_id)
    temporary = PPTX.with_suffix(".pptx.tmp")
    with warnings.catch_warnings():
        warnings.filterwarnings("ignore", message="Duplicate name: .*", category=UserWarning)
        presentation.save(temporary)
    normalize_unique_zip(temporary)
    temporary.replace(PPTX)


def fix_workbook() -> None:
    workbook = load_workbook(XLSX)
    sheet = workbook["Dashboard"]
    for chart in list(sheet._charts):
        sheet._charts.remove(chart)
    chart = BarChart()
    chart.type = "col"
    chart.style = 10
    chart.title = "Distribución comparada"
    chart.y_axis.title = "Resultado"
    chart.x_axis.title = "Categoría"
    data = Reference(sheet, min_col=2, max_col=3, min_row=8, max_row=13)
    categories = Reference(sheet, min_col=1, min_row=9, max_row=13)
    chart.add_data(data, titles_from_data=True)
    chart.set_categories(categories)
    chart.height = 8.5
    chart.width = 15
    chart.legend.position = "b"
    chart.dLbls = DataLabelList()
    chart.dLbls.showVal = True
    colors = ["4C2B46", "E97700"]
    for series, color in zip(chart.series, colors):
        series.graphicalProperties.solidFill = color
        series.graphicalProperties.line.solidFill = color
    sheet.add_chart(chart, "E8")
    if "Galería visual" in workbook.sheetnames:
        del workbook["Galería visual"]
    gallery = workbook.create_sheet("Galería visual")
    gallery.sheet_view.showGridLines = False
    gallery.append(["Categoría", "Resultado", "Referencia"])
    for row in [("A", 82, 70), ("B", 64, 58), ("C", 47, 52), ("D", 31, 40), ("E", 18, 25)]:
        gallery.append(row)
    for cell in gallery[1]:
        cell.font = cell.font.copy(name="Aptos", bold=True, color="FFFFFF")
        cell.fill = cell.fill.copy(fill_type="solid", fgColor="4C2B46")
    data = Reference(gallery, min_col=2, max_col=3, min_row=1, max_row=6)
    categories = Reference(gallery, min_col=1, min_row=2, max_row=6)
    chart_specs = []
    for title, kind, grouping in (
        ("Columnas agrupadas", "col", "clustered"),
        ("Columnas apiladas", "col", "stacked"),
        ("Columnas 100 %", "col", "percentStacked"),
        ("Barras agrupadas", "bar", "clustered"),
        ("Barras apiladas", "bar", "stacked"),
    ):
        item = BarChart(); item.type = kind; item.grouping = grouping; item.title = title; chart_specs.append(item)
    pie = PieChart(); pie.title = "Gráfico de torta"; chart_specs.append(pie)
    doughnut = DoughnutChart(); doughnut.title = "Gráfico de anillo"; chart_specs.append(doughnut)
    line = LineChart(); line.title = "Evolución temporal"; chart_specs.append(line)
    area = AreaChart(); area.title = "Área comparada"; chart_specs.append(area)
    radar = RadarChart(); radar.title = "Perfil comparado"; chart_specs.append(radar)
    anchors = ["E2", "M2", "E18", "M18", "E34", "M34", "E50", "M50", "E66", "M66"]
    for item, anchor in zip(chart_specs, anchors):
        item.add_data(data, titles_from_data=True)
        item.set_categories(categories)
        item.height = 7.5; item.width = 13.5
        item.legend.position = "b"
        for series, color in zip(item.series, ("4C2B46", "E97700")):
            series.graphicalProperties.solidFill = color
            series.graphicalProperties.line.solidFill = color
        gallery.add_chart(item, anchor)
    gallery.freeze_panes = "A2"
    gallery.column_dimensions["A"].width = 18
    gallery.column_dimensions["B"].width = 14
    gallery.column_dimensions["C"].width = 14
    temporary = XLSX.with_suffix(".xlsx.tmp")
    workbook.save(temporary)
    temporary.replace(XLSX)


def main() -> int:
    build_oficio()
    build_circular()
    build_agenda()
    build_acta()
    build_minuta()
    build_programa()
    build_invitation()
    report_charts = [CHART_LIBRARY[index] for index in (0, 5, 7)]
    for name in ("STGND_Informe_Maestro_v9.docx", "STGND_Informe_Fondo_Color_Maestro_v9.docx"):
        inject_native_charts(OUT / name, report_charts)
    graphics_path = OUT / "STGND_Documento_Graficos_Maestro_v9.docx"
    inject_native_charts(graphics_path, CHART_LIBRARY)
    add_structured_data_examples(graphics_path)
    formal = {"STGND_Oficio_Maestro_v9.docx", "STGND_Oficio_Circular_Maestro_v9.docx"}
    for path in OUT.glob("*.docx"):
        if path.name not in formal:
            apply_semantic_styles(path)
    for name in ("STGND_Oficio_Maestro_v9.docx", "STGND_Oficio_Circular_Maestro_v9.docx", "STGND_Acta_Reunion_Maestro_v9.docx"):
        normalize_signature_box(OUT / name, ensure=True)
    for name in formal:
        enforce_formal_black(OUT / name)
    fix_presentation()
    fix_workbook()
    print("OK: formatos institucionales simples y gráficos Office nativos aplicados")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
